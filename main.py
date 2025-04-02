import sys
import re
from PyQt5.QtWidgets import (QApplication, QMainWindow, QLabel, QWidget, 
                             QVBoxLayout, QFileDialog, QPushButton, QListWidget,
                             QLineEdit, QMessageBox)
from PyQt5.QtCore import Qt
from pptx import Presentation
from docx import Document

class PowerpointFile():
    
    def __init__(self, filePath):
        self.filePath = filePath
        self.fileName = filePath.split("/")[-1]

class Line():
    
    def __init__(self, text, bold=False, underline=False, startOfPresentation=False, startOfSlide=False):
        self.text = text
        self.bold = bold
        self.underline = underline
        self.startOfPresentation = startOfPresentation
        self.startOfSlide = startOfSlide

class MainWindow(QMainWindow):
    
    selectedPowerpointFiles = []
    lines = []
    selectedFolder = None
    
    def __init__(self):
        super().__init__()
        self.initUI()
        self.setWindowTitle("Powerpoint to Word")
        
    def initUI(self):
        centralWidget = QWidget()
        self.setCentralWidget(centralWidget)
        
        heading = QLabel("Powerpoint to Word")
        listLabel = QLabel("List of Powerpoint Files")
        self.listWidget = QListWidget()
        selectButton = QPushButton("Select Files")
        folderButton = QPushButton("Select Folder")
        inputLabel = QLabel("Output File Name")
        self.nameInput = QLineEdit("Output")
        executeButton = QPushButton("Execute")
        
        selectButton.clicked.connect(self.selectFiles)
        folderButton.clicked.connect(self.selectFolder)
        executeButton.clicked.connect(self.execute)
        
        heading.setStyleSheet("font-size: 20px; font-weight: bold; padding: 10px")
        inputLabel.setAlignment(Qt.AlignCenter)
        heading.setAlignment(Qt.AlignCenter)
        
        vbox = QVBoxLayout()
        vbox.addWidget(heading)
        vbox.addWidget(listLabel)
        vbox.addWidget(self.listWidget)
        vbox.addWidget(selectButton)
        vbox.addWidget(folderButton)
        vbox.addWidget(inputLabel)
        vbox.addWidget(self.nameInput)
        vbox.addWidget(executeButton)
        
        centralWidget.setLayout(vbox)
        
    def selectFiles(self):
        filePaths, _ = QFileDialog.getOpenFileNames(self, "Select Powerpoint Files", "", "Powerpoint Files (*.pptx)")
        self.selectedPowerpointFiles = [PowerpointFile(filePath) for filePath in filePaths]
        self.listWidget.clear()
        for powerpointFile in self.selectedPowerpointFiles:
            self.listWidget.addItem(powerpointFile.fileName)
            
    def selectFolder(self):
        self.selectedFolder = QFileDialog.getExistingDirectory(self, "Select Folder")
            
    def execute(self):
        document = Document()
        
        for powerpointFile in self.selectedPowerpointFiles:
            self.exportLinesFromPresentation(powerpointFile)
            self.addToWordDocument(document)  
            self.lines.clear()
        
        document.save(f"{self.selectedFolder}/{self.nameInput.text()}.docx")
        self.selectedPowerpointFiles.clear()
        self.listWidget.clear()
        okMsg = QMessageBox()
        okMsg.setWindowTitle("Finished")
        okMsg.setText("Word Document Was Created")
        okMsg.setStandardButtons(QMessageBox.StandardButton.Ok)
        okMsg.exec_()
            
    def exportLinesFromPresentation(self, powerpointFile):
        presentation = Presentation(powerpointFile.filePath)

        for slide in presentation.slides:
            startOfSlide = True
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        isBold = False
                        isUnderline = False
                        for run in paragraph.runs:
                            if run.font.bold:
                                isBold = True
                            if run.font.underline:
                                isUnderline = True
                        
                        
                        isStartOfPresentation = len(self.lines) == 0
                        line = Line(paragraph.text, isBold, isUnderline, isStartOfPresentation, startOfSlide)
                        self.lines.append(line)
                        startOfSlide = False
                    
    def addToWordDocument(self, document):
        for line in self.lines:
            cleanedText = self.cleanText(line.text)
            
            if line.startOfPresentation:
                document.add_heading(cleanedText, level=1)
            elif line.startOfSlide:
                document.add_heading(cleanedText, level=3)
            else:
                paragraph = document.add_paragraph()
                run = paragraph.add_run(cleanedText)
                run.bold = line.bold
                run.underline = line.underline
        
    def cleanText(self, text):
        # Makes text compatible with XML
        cleanedText = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', text)
        return cleanedText
    

def loadStylesheet(filename):
    with open(filename, "r") as file:
        return file.read()    
    
def main():
    app = QApplication(sys.argv)
    window = MainWindow()
    stylesheet = loadStylesheet("stylesheets.css")
    window.setStyleSheet(stylesheet)
    window.show()
    sys.exit(app.exec_())

if __name__ == '__main__':
    main()